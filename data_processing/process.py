import os
import fitz  
import docx2txt
from pptx import Presentation
import tiktoken
import uuid
import re
import shutil
import json  

IMAGE_DIR = "../outputs/images"
CHUNK_DIR = "../outputs/chunks"
METADATA_DIR = "../outputs/metadata"  
PARENT_CHUNK_TOKENS = 1024
CHILD_CHUNK_TOKENS = 256

tokenizer = tiktoken.get_encoding("cl100k_base")

def create_dirs():
    """Create output directories if they don't exist."""
    os.makedirs(IMAGE_DIR, exist_ok=True)
    os.makedirs(CHUNK_DIR, exist_ok=True)
    os.makedirs(METADATA_DIR, exist_ok=True)  

def preprocess_text(text):
    """Basic text cleaning."""
    text = text.lower()
    text = re.sub(r'\s+', ' ', text)
    return text.strip()


def extract_from_pdf(file_path, doc_id):
    """Extracts text and images from a PDF file."""
    text_content = ""
    image_files = []  
    doc = fitz.open(file_path)
    for page_num, page in enumerate(doc):
        text_content += page.get_text() + "\n"
        for img_index, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]
            
            image_filename = f"{doc_id}_page{page_num+1}_img{img_index+1}.{image_ext}"
            image_save_path = os.path.join(IMAGE_DIR, image_filename)
            with open(image_save_path, "wb") as img_file:
                img_file.write(image_bytes)
            
            text_content += f"[Image: {image_filename}]\n"
            image_files.append(image_filename)  
            
    return text_content, image_files  

def extract_from_docx(file_path, doc_id):
    """Extracts text and images from a DOCX file using docx2txt."""
    image_files = []  
    image_subfolder_rel = f"{doc_id}_images"
    image_subfolder_abs = os.path.join(IMAGE_DIR, image_subfolder_rel)
    text_content = docx2txt.process(file_path, image_subfolder_abs)
    
    if os.path.exists(image_subfolder_abs):
        for img_name in os.listdir(image_subfolder_abs):
            image_rel_path = f"{image_subfolder_rel}/{img_name}" 
            text_content += f"\n[Image: {image_rel_path}]"
            image_files.append(image_rel_path)  
            
    return text_content, image_files  

def extract_from_pptx(file_path, doc_id):
    """Extracts text and images from a PPTX file."""
    text_content = ""
    image_files = []  
    pres = Presentation(file_path)
    for slide_num, slide in enumerate(pres.slides):
        text_content += f"\n--- Slide {slide_num+1} ---\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content += shape.text + "\n"
            if "picture" in str(type(shape)).lower():
                image_bytes = shape.image.blob
                image_ext = shape.image.ext
                image_filename = f"{doc_id}_slide{slide_num+1}_{shape.shape_id}.{image_ext}"
                image_save_path = os.path.join(IMAGE_DIR, image_filename)
                with open(image_save_path, "wb") as img_file:
                    img_file.write(image_bytes)
                text_content += f"[Image: {image_filename}]\n"
                image_files.append(image_filename)  
                
    return text_content, image_files  

def extract_from_text(file_path):
    """Extracts text from a TXT file."""
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read(), []  

def extract_from_markdown(file_path, doc_id):
    """
    Extracts text from a Markdown file, finds linked images,
    copies them to the image directory, and replaces the links with placeholders.
    """
    base_dir = os.path.dirname(file_path)
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()

    img_pattern = re.compile(r'!\[(.*?)\]\((.*?)\)')
    image_files = []  
    
    matches = list(img_pattern.finditer(content))
    img_index = 0

    for match in reversed(matches):
        alt_text = match.group(1)
        original_path = match.group(2)
        
        source_image_path = os.path.join(base_dir, original_path)
        
        if os.path.exists(source_image_path):
            img_index += 1
            ext = os.path.splitext(source_image_path)[1]
            
            new_image_filename = f"{doc_id}_md_img_{img_index}{ext}"
            dest_image_path = os.path.join(IMAGE_DIR, new_image_filename)
            
            shutil.copy(source_image_path, dest_image_path)
            
            placeholder = f"[Image: {new_image_filename} | Alt Text: {alt_text}]"
            content = content[:match.start()] + placeholder + content[match.end():]
            image_files.append(new_image_filename)  
        else:
            print(f"Warning: Image not found at path: {source_image_path}")

    return content, image_files  

def chunk_text(text, doc_id):
    """
    Performs content-aware hierarchical chunking.
    Returns lists of the chunk filenames created.
    """
    paragraphs = text.split('\n')
    paragraphs = [p.strip() for p in paragraphs if p.strip()]

    child_chunks = []
    for para in paragraphs:
        if len(tokenizer.encode(para)) > CHILD_CHUNK_TOKENS:
            sentences = re.split(r'(?<=[.!?])\s+', para)
            current_chunk = ""
            for sentence in sentences:
                if len(tokenizer.encode(current_chunk + sentence)) <= CHILD_CHUNK_TOKENS:
                    current_chunk += " " + sentence
                else:
                    child_chunks.append(current_chunk.strip())
                    current_chunk = sentence
            child_chunks.append(current_chunk.strip())
        else:
            child_chunks.append(para)

    parent_chunks = []
    current_parent = ""
    for chunk in child_chunks:
        if len(tokenizer.encode(current_parent)) + len(tokenizer.encode(chunk)) <= PARENT_CHUNK_TOKENS:
            current_parent += "\n" + chunk
        else:
            parent_chunks.append(current_parent.strip())
            current_parent = chunk
    parent_chunks.append(current_parent.strip())
    
    child_chunk_files = []
    parent_chunk_files = []
    
    for i, chunk in enumerate(child_chunks):
        chunk_filename = f"{doc_id}_child_chunk_{i+1}.txt"
        with open(os.path.join(CHUNK_DIR, chunk_filename), 'w', encoding='utf-8') as f:
            f.write(chunk)
        child_chunk_files.append(chunk_filename)
            
    for i, chunk in enumerate(parent_chunks):
        chunk_filename = f"{doc_id}_parent_chunk_{i+1}.txt"
        with open(os.path.join(CHUNK_DIR, chunk_filename), 'w', encoding='utf-8') as f:
            f.write(chunk)
        parent_chunk_files.append(chunk_filename)
            
    print(f"Created {len(child_chunk_files)} child chunks and {len(parent_chunk_files)} parent chunks.")
    return child_chunk_files, parent_chunk_files  

def process_file(file_path):
    """
    Main function to process a single file.
    Detects file type, extracts content, chunks it, and saves a metadata manifest.
    """
    create_dirs()
    
    doc_id = os.path.basename(file_path).split('.')[0]
    print(f"Processing {file_path} with ID: {doc_id}")
    
    ext = os.path.splitext(file_path)[1].lower()
    full_text = ""
    image_filenames = []  
    
    if ext == '.pdf':
        full_text, image_filenames = extract_from_pdf(file_path, doc_id)
    elif ext == '.docx':
        full_text, image_filenames = extract_from_docx(file_path, doc_id)
    elif ext == '.pptx':
        full_text, image_filenames = extract_from_pptx(file_path, doc_id)
    elif ext == '.md':
        full_text, image_filenames = extract_from_markdown(file_path, doc_id)
    elif ext == '.txt':
        full_text, image_filenames = extract_from_text(file_path) 
    else:
        print(f"Unsupported file type: {ext}")
        return

    cleaned_text = preprocess_text(full_text)
    child_chunk_files, parent_chunk_files = chunk_text(cleaned_text, doc_id)
    
    metadata = {
        "doc_id": doc_id,
        "source_file": os.path.abspath(file_path),
        "assets": {
            "images": image_filenames,
            "child_chunks": child_chunk_files,
            "parent_chunks": parent_chunk_files
        }
    }
    
    metadata_filename = f"{doc_id}.json"
    metadata_save_path = os.path.join(METADATA_DIR, metadata_filename)
    
    with open(metadata_save_path, 'w', encoding='utf-8') as f:
        json.dump(metadata, f, indent=4)
        
    print(f"Processing complete. Metadata manifest saved to {metadata_save_path}")


if __name__ == '__main__':
    import sys
    if len(sys.argv) < 2:
        print("Usage: python process_documents.py <path_to_file>")
    else:
        file_to_process = sys.argv[1]
        if not os.path.exists(file_to_process):
            print(f"Error: File not found at {file_to_process}")
        else:
            process_file(file_to_process)